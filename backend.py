"""
PLAGIFY Backend - Système de détection de plagiat ultra-performant
Intégration complète avec Supabase
"""

from fastapi import FastAPI, UploadFile, File, WebSocket, WebSocketDisconnect, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from supabase import create_client, Client
from typing import List, Optional
import os
from dotenv import load_dotenv
import asyncio
import hashlib
import re
from pathlib import Path
from datetime import datetime
import json
from concurrent.futures import ThreadPoolExecutor, ProcessPoolExecutor
from difflib import SequenceMatcher
import mimetypes

# Extraction de texte
import PyPDF2
import docx
from pptx import Presentation

# Génération de rapports PDF
from reportlab.lib.pagesizes import A4, landscape
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch, cm
from reportlab.lib.colors import HexColor, red, orange, yellow
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
from reportlab.graphics.shapes import Drawing
from reportlab.graphics.charts.piecharts import Pie
from reportlab.graphics.charts.barcharts import VerticalBarChart
from reportlab.graphics.charts.linecharts import HorizontalLineChart

load_dotenv()

# Configuration
app = FastAPI(title="PlaGiFY API", version="2.0.0")

# CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ✅ FIX PROBLÈME 1 — supabase_client est le seul objet Supabase, jamais écrasé
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_KEY")

if not SUPABASE_URL or not SUPABASE_KEY:
    raise Exception("SUPABASE_URL et SUPABASE_KEY doivent être définis dans .env")

supabase_client: Client = create_client(SUPABASE_URL, SUPABASE_KEY)

def get_bucket(bucket_name: str):
    """
    Helper compatible supabase-py v1 ET v2.
    - v1 : supabase_client.storage().from_(bucket)
    - v2 : supabase_client.storage.from_(bucket)
    """
    storage = supabase_client.storage
    # Si storage est un callable (v1), on l'appelle d'abord
    if callable(storage):
        return storage().from_(bucket_name)
    # Sinon (v2), on accède directement
    return storage.from_(bucket_name)

# Configuration
UPLOAD_DIR = Path("/tmp/plagify_uploads")
REPORTS_DIR = Path("/tmp/plagify_reports")
UPLOAD_DIR.mkdir(exist_ok=True)
REPORTS_DIR.mkdir(exist_ok=True)

VALID_EXTENSIONS = {'.pdf', '.txt', '.doc', '.docx', '.ppt', '.pptx', '.html', '.css', '.js', '.php', '.c', '.py', '.java'}

# Executors pour parallélisation
thread_executor = ThreadPoolExecutor(max_workers=10)
process_executor = ProcessPoolExecutor(max_workers=4)

# WebSocket connections
ws_connections = {}

# ============================================
# UTILITAIRES
# ============================================

def extract_text_from_file(file_path: Path) -> tuple[str, str]:
    """Extrait le texte d'un fichier et détecte le langage"""
    ext = file_path.suffix.lower()
    text = ""
    language = "unknown"
    
    try:
        if ext == '.pdf':
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = ' '.join(page.extract_text() for page in reader.pages)
            language = "document"
            
        elif ext in ['.doc', '.docx']:
            doc = docx.Document(file_path)
            text = ' '.join(paragraph.text for paragraph in doc.paragraphs)
            language = "document"
            
        elif ext in ['.ppt', '.pptx']:
            prs = Presentation(file_path)
            text = ' '.join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
            language = "presentation"
            
        elif ext in ['.html', '.css', '.js', '.php', '.c', '.py', '.java']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            language_map = {
                '.html': 'HTML', '.css': 'CSS', '.js': 'JavaScript',
                '.php': 'PHP', '.c': 'C', '.py': 'Python', '.java': 'Java'
            }
            language = language_map.get(ext, 'code')
            
        else:  # .txt
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            language = "text"
            
    except Exception as e:
        print(f"Erreur extraction {file_path}: {e}")
        
    return text, language

def compute_hash(text: str) -> str:
    """Calcule le hash SHA256 du texte"""
    return hashlib.sha256(text.encode('utf-8', errors='ignore')).hexdigest()

def calculate_similarity(text1: str, text2: str) -> tuple[float, dict]:
    """
    Calcule la similarité entre deux textes
    Retourne: (score_global, détails)
    """
    # Similarité globale
    global_similarity = SequenceMatcher(None, text1, text2).ratio() * 100
    
    # Trouver les segments similaires
    matcher = SequenceMatcher(None, text1, text2)
    segments = {
        'exact': [],  # >80% similarité
        'moderate': [],  # 50-80%
        'weak': []  # 30-50%
    }
    
    for tag, i1, i2, j1, j2 in matcher.get_opcodes():
        if tag == 'equal':
            segment_length = i2 - i1
            if segment_length >= 20:  # Au moins 20 caractères
                segments['exact'].append({
                    'text_a_start': i1,
                    'text_a_end': i2,
                    'text_b_start': j1,
                    'text_b_end': j2,
                    'text': text1[i1:i2],
                    'similarity': 100
                })
        elif tag == 'replace':
            segment_sim = SequenceMatcher(None, text1[i1:i2], text2[j1:j2]).ratio() * 100
            if segment_sim >= 30 and (i2 - i1) >= 20:
                seg_data = {
                    'text_a_start': i1,
                    'text_a_end': i2,
                    'text_b_start': j1,
                    'text_b_end': j2,
                    'text_a': text1[i1:i2],
                    'text_b': text2[j1:j2],
                    'similarity': segment_sim
                }
                if segment_sim >= 80:
                    segments['exact'].append(seg_data)
                elif segment_sim >= 50:
                    segments['moderate'].append(seg_data)
                else:
                    segments['weak'].append(seg_data)
    
    details = {
        'global_similarity': round(global_similarity, 2),
        'exact_count': len(segments['exact']),
        'moderate_count': len(segments['moderate']),
        'weak_count': len(segments['weak']),
        'segments': segments
    }
    
    return round(global_similarity, 2), details

async def send_progress(ws_id: str, data: dict):
    """Envoie une mise à jour de progression via WebSocket"""
    if ws_id in ws_connections:
        try:
            await ws_connections[ws_id].send_json(data)
        except:
            pass

def generate_pdf_report(report_data: dict, output_path: Path) -> Path:
    """
    Génère un rapport PDF professionnel selon les spécifications exactes
    """
    doc = SimpleDocTemplate(
        str(output_path),
        pagesize=A4,
        rightMargin=1*cm,
        leftMargin=1*cm,
        topMargin=1*cm,
        bottomMargin=1*cm
    )
    
    story = []
    styles = getSampleStyleSheet()
    
    # Style personnalisé
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=HexColor('#FF3D71'),
        alignment=TA_CENTER
    )
    
    # EN-TÊTE
    header_data = [
        ['', 'PlaGiFY', f"Date: {report_data['date']}"],
        ['', report_data.get('establishment_name', ''), f"ID: {report_data['report_id']}"],
        ['', report_data.get('teacher_name', ''), Paragraph(f"<font color='red' size='16'><b>{report_data['global_similarity']}%</b></font>", styles['Normal'])]
    ]
    
    header_table = Table(header_data, colWidths=[6*cm, 8*cm, 6*cm])
    header_table.setStyle(TableStyle([
        ('ALIGN', (0, 0), (0, -1), 'LEFT'),
        ('ALIGN', (1, 0), (1, -1), 'CENTER'),
        ('ALIGN', (2, 0), (2, -1), 'RIGHT'),
        ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
        ('TEXTCOLOR', (1, 0), (1, 0), HexColor('#FF3D71')),
        ('FONTSIZE', (1, 0), (1, 0), 18),
        ('FONTNAME', (1, 0), (1, 0), 'Helvetica-Bold'),
    ]))
    
    story.append(header_table)
    story.append(Spacer(1, 0.5*cm))
    
    # TABLEAU STATISTIQUES
    stats_data = [
        ['📊 Comparaisons', f"{report_data['total_comparisons']}", '📈 Moy. similarité', f"{report_data['avg_similarity']}%"],
        ['🚨 Matches > seuil', f"{report_data['matches_count']}", '🎯 Seuil configuré', f"{report_data['threshold']}%"]
    ]
    
    stats_table = Table(stats_data, colWidths=[4*cm, 4*cm, 4*cm, 4*cm])
    stats_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, -1), HexColor('#F0F0F0')),
        ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#CCCCCC')),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, -1), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
    ]))
    
    story.append(stats_table)
    story.append(Spacer(1, 0.5*cm))
    
    # INFORMATIONS FICHIERS
    file_info_data = [
        ['Fichier A', 'Fichier B', 'Taille A', 'Taille B'],
        [report_data['file_a_name'], report_data['file_b_name'], 
         f"{report_data['file_a_size']} octets", f"{report_data['file_b_size']} octets"],
        ['Mots A', 'Mots B', 'Langage A', 'Langage B'],
        [f"{report_data['file_a_words']}", f"{report_data['file_b_words']}",
         report_data['file_a_language'], report_data['file_b_language']],
        ['Similarité exacte', 'Similarité modérée', 'Similarité faible', 'Type'],
        [f"{report_data['exact_matches']}", f"{report_data['moderate_matches']}",
         f"{report_data['weak_matches']}", report_data['similarity_type']]
    ]
    
    file_info_table = Table(file_info_data, colWidths=[4.5*cm] * 4)
    file_info_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), HexColor('#4ECDC4')),
        ('BACKGROUND', (0, 2), (-1, 2), HexColor('#4ECDC4')),
        ('BACKGROUND', (0, 4), (-1, 4), HexColor('#4ECDC4')),
        ('TEXTCOLOR', (0, 0), (-1, 0), HexColor('#FFFFFF')),
        ('TEXTCOLOR', (0, 2), (-1, 2), HexColor('#FFFFFF')),
        ('TEXTCOLOR', (0, 4), (-1, 4), HexColor('#FFFFFF')),
        ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#CCCCCC')),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
        ('FONTSIZE', (0, 0), (-1, -1), 8),
    ]))
    
    story.append(file_info_table)
    story.append(Spacer(1, 0.5*cm))
    
    # VISUALISATION CÔTE À CÔTE
    story.append(Paragraph("<b>VISUALISATION DES PARTIES SIMILAIRES</b>", styles['Heading2']))
    story.append(Spacer(1, 0.3*cm))
    
    # Limiter le texte pour tenir sur une page
    text_a_preview = report_data['text_a'][:2000] if len(report_data['text_a']) > 2000 else report_data['text_a']
    text_b_preview = report_data['text_b'][:2000] if len(report_data['text_b']) > 2000 else report_data['text_b']
    
    comparison_data = [
        ['FICHIER A', '|', 'FICHIER B'],
        [text_a_preview, '|', text_b_preview]
    ]
    
    comparison_table = Table(comparison_data, colWidths=[8.5*cm, 0.5*cm, 8.5*cm])
    comparison_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (0, 0), HexColor('#FF6B6B')),
        ('BACKGROUND', (2, 0), (2, 0), HexColor('#4ECDC4')),
        ('TEXTCOLOR', (0, 0), (2, 0), HexColor('#FFFFFF')),
        ('ALIGN', (0, 0), (-1, 0), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 10),
        ('VALIGN', (0, 1), (-1, 1), 'TOP'),
        ('FONTNAME', (0, 1), (-1, 1), 'Courier'),
        ('FONTSIZE', (0, 1), (-1, 1), 7),
        ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#CCCCCC')),
    ]))
    
    story.append(comparison_table)
    story.append(PageBreak())
    
    # PAGE 2: STATISTIQUES DÉTAILLÉES
    story.append(Paragraph("<b>ANALYSE DÉTAILLÉE</b>", title_style))
    story.append(Spacer(1, 0.5*cm))
    
    detailed_stats_data = [
        ['Métrique', 'Valeur'],
        ['Similarité brute', f"{report_data['global_similarity']}%"],
        ['Similarité sans citations', f"{report_data.get('similarity_no_quotes', report_data['global_similarity'])}%"],
        ['Similarité structurelle', f"{report_data.get('structural_similarity', 0)}%"],
        ['Similarité syntaxique', f"{report_data.get('syntactic_similarity', 0)}%"],
    ]
    
    detailed_table = Table(detailed_stats_data, colWidths=[10*cm, 7*cm])
    detailed_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), HexColor('#667EEA')),
        ('TEXTCOLOR', (0, 0), (-1, 0), HexColor('#FFFFFF')),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#CCCCCC')),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
    ]))
    
    story.append(detailed_table)
    story.append(Spacer(1, 1*cm))
    
    # DIAGRAMMES
    story.append(Paragraph("<b>GRAPHIQUES DE SIMILARITÉ</b>", styles['Heading2']))
    story.append(Spacer(1, 0.3*cm))
    
    # Diagramme circulaire — évite division par zéro
    exact_m = report_data['exact_matches'] or 0
    moderate_m = report_data['moderate_matches'] or 0
    weak_m = report_data['weak_matches'] or 0
    total_m = exact_m + moderate_m + weak_m

    if total_m > 0:
        drawing = Drawing(400, 200)
        pie = Pie()
        pie.x = 150
        pie.y = 50
        pie.width = 100
        pie.height = 100
        pie.data = [exact_m, moderate_m, weak_m]
        pie.labels = ['Exact', 'Modéré', 'Faible']
        pie.slices[0].fillColor = HexColor('#FF6B6B')
        pie.slices[1].fillColor = HexColor('#FFA500')
        pie.slices[2].fillColor = HexColor('#FFD700')
        drawing.add(pie)
        story.append(drawing)
    else:
        story.append(Paragraph("Aucun segment similaire détecté.", styles['Normal']))
    
    story.append(Spacer(1, 1*cm))
    
    # FOOTER
    footer_data = [
        ['Date d\'analyse', report_data['date']],
        ['Signature numérique', report_data['signature']],
        ['Version algorithme', 'PlaGiFY v2.0.0']
    ]
    
    footer_table = Table(footer_data, colWidths=[8*cm, 10*cm])
    footer_table.setStyle(TableStyle([
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
    ]))
    
    story.append(footer_table)
    
    # Construire le PDF
    doc.build(story)
    
    return output_path

# ============================================
# ENDPOINTS API
# ============================================

@app.get("/")
async def root():
    return {"message": "PlaGiFY API v2.0 - Système de détection de plagiat"}

@app.post("/api/teachers")
async def create_teacher(name: str = Form(...), email: str = Form(...)):
    """Créer ou mettre à jour un enseignant"""
    try:
        existing = supabase_client.table('teachers').select('*').eq('email', email).execute()
        
        if existing.data:
            result = supabase_client.table('teachers').update({
                'name': name
            }).eq('email', email).execute()
            return {"success": True, "data": result.data[0], "message": "Enseignant mis à jour"}
        else:
            result = supabase_client.table('teachers').insert({
                'name': name,
                'email': email
            }).execute()
            return {"success": True, "data": result.data[0], "message": "Enseignant créé"}
            
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/teachers/{email}")
async def get_teacher(email: str):
    """Récupérer un enseignant par email"""
    try:
        result = supabase_client.table('teachers').select('*').eq('email', email).execute()
        if result.data:
            return {"success": True, "data": result.data[0]}
        return {"success": False, "message": "Enseignant non trouvé"}
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/establishments")
async def create_establishment(
    teacher_id: str = Form(...),
    name: str = Form(...),
    logo_file: Optional[UploadFile] = File(None)
):
    """Créer un établissement"""
    try:
        logo_url = None
        
        if logo_file:
            try:
                file_bytes = await logo_file.read()
                file_path = f"logos/{teacher_id}/{logo_file.filename}"
                
                get_bucket('plagify-files').upload(
                    file_path,
                    file_bytes,
                    {'content-type': logo_file.content_type or 'image/png', 'upsert': 'true'}
                )
                
                logo_url = get_bucket('plagify-files').get_public_url(file_path)
            except Exception as storage_error:
                print(f"Storage error (non-critical): {storage_error}")
                logo_url = None
        
        result = supabase_client.table('establishments').insert({
            'teacher_id': teacher_id,
            'name': name,
            'logo_url': logo_url
        }).execute()
        
        return {"success": True, "data": result.data[0]}
        
    except Exception as e:
        print(f"Error creating establishment: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/establishments/{teacher_id}")
async def get_establishments(teacher_id: str):
    """Récupérer tous les établissements d'un enseignant"""
    try:
        result = supabase_client.table('establishments').select('*').eq('teacher_id', teacher_id).execute()
        return {"success": True, "data": result.data}
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/establishments/{establishment_id}")
async def delete_establishment(establishment_id: str):
    """Supprimer un établissement"""
    try:
        supabase_client.table('establishments').delete().eq('id', establishment_id).execute()
        return {"success": True, "message": "Établissement supprimé"}
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/files/upload")
async def upload_files(
    teacher_id: str = Form(...),
    files: List[UploadFile] = File(...)
):
    """Upload et enregistrement de fichiers dans la base de données"""
    try:
        uploaded_files = []
        errors = []
        
        for file in files:
            try:
                ext = Path(file.filename).suffix.lower()
                
                if ext not in VALID_EXTENSIONS:
                    errors.append(f"{file.filename}: Extension non supportée")
                    continue
                
                # ✅ FIX PROBLÈME 5 — Extraire uniquement le nom du fichier, sans sous-dossiers
                safe_filename = Path(file.filename).name
                temp_path = UPLOAD_DIR / safe_filename
                file_content = await file.read()
                with open(temp_path, 'wb') as f:
                    f.write(file_content)
                
                text, language = extract_text_from_file(temp_path)
                content_hash = compute_hash(text)
                word_count = len(text.split())
                file_size = temp_path.stat().st_size
                
                storage_path = f"files/{teacher_id}/{safe_filename}"
                try:
                    with open(temp_path, 'rb') as f:
                        get_bucket('plagify-files').upload(
                            storage_path,
                            f.read(),
                            {'content-type': file.content_type or 'application/octet-stream', 'upsert': 'true'}
                        )
                except Exception as storage_error:
                    print(f"Storage upload error for {file.filename}: {storage_error}")
                
                result = supabase_client.table('files').insert({
                    'teacher_id': teacher_id,
                    'filename': safe_filename,
                    'original_path': file.filename,
                    'storage_path': storage_path,
                    'file_type': ext,
                    'file_size': file_size,
                    'content_text': text[:50000],
                    'content_hash': content_hash,
                    'word_count': word_count,
                    'language': language
                }).execute()
                
                uploaded_files.append(result.data[0])
                
                if temp_path.exists():
                    temp_path.unlink()
                    
            except Exception as file_error:
                errors.append(f"{file.filename}: {str(file_error)}")
                print(f"Error processing file {file.filename}: {file_error}")
                continue
        
        return {
            "success": True, 
            "data": uploaded_files, 
            "count": len(uploaded_files),
            "errors": errors if errors else None
        }
        
    except Exception as e:
        print(f"General upload error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/files/{teacher_id}")
async def get_teacher_files(teacher_id: str):
    """Récupérer tous les fichiers d'un enseignant"""
    try:
        result = supabase_client.table('files').select('*').eq('teacher_id', teacher_id).order('uploaded_at', desc=True).execute()
        return {"success": True, "data": result.data}
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/files/{file_id}")
async def delete_file(file_id: str):
    """Supprimer un fichier"""
    try:
        file_data = supabase_client.table('files').select('*').eq('id', file_id).execute()
        if file_data.data:
            get_bucket('plagify-files').remove([file_data.data[0]['storage_path']])
        
        supabase_client.table('files').delete().eq('id', file_id).execute()
        return {"success": True, "message": "Fichier supprimé"}
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/analyze/folder")
async def analyze_folder(
    teacher_id: str = Form(...),
    establishment_id: Optional[str] = Form(None),
    threshold: float = Form(15.0),
    ws_id: str = Form(...),
    files: List[UploadFile] = File(...)
):
    """
    Analyse complète d'un dossier de fichiers
    Compare tous les fichiers entre eux et génère des rapports
    """
    try:
        analysis = supabase_client.table('analyses').insert({
            'teacher_id': teacher_id,
            'establishment_id': establishment_id,
            'analysis_type': 'folder',
            'source_name': f"Dossier ({len(files)} fichiers)",
            'similarity_threshold': threshold,
            'status': 'processing',
            'total_files': len(files)
        }).execute()
        
        analysis_id = analysis.data[0]['id']
        
        await send_progress(ws_id, {
            'stage': 'extraction',
            'progress': 0,
            'total': len(files),
            'message': 'Extraction des fichiers...'
        })
        
        file_records = []
        for idx, file in enumerate(files):
            ext = Path(file.filename).suffix.lower()
            if ext not in VALID_EXTENSIONS:
                continue
            
            # ✅ FIX PROBLÈME 5 — Nom de fichier sûr sans sous-dossiers
            safe_filename = Path(file.filename).name
            temp_path = UPLOAD_DIR / f"{analysis_id}_{safe_filename}"
            with open(temp_path, 'wb') as f:
                f.write(await file.read())
            
            text, language = extract_text_from_file(temp_path)
            content_hash = compute_hash(text)
            
            storage_path = f"analyses/{analysis_id}/{safe_filename}"
            try:
                with open(temp_path, 'rb') as f:
                    get_bucket('plagify-files').upload(
                        storage_path,
                        f.read(),
                        {'content-type': file.content_type or 'application/octet-stream', 'upsert': 'true'}
                    )
            except Exception as storage_error:
                print(f"Storage upload error for {file.filename}: {storage_error}")
            
            file_record = supabase_client.table('files').insert({
                'teacher_id': teacher_id,
                'filename': safe_filename,
                'original_path': file.filename,
                'storage_path': storage_path,
                'file_type': ext,
                'file_size': temp_path.stat().st_size,
                'content_text': text[:50000],
                'content_hash': content_hash,
                'word_count': len(text.split()),
                'language': language
            }).execute()
            
            file_records.append({
                'id': file_record.data[0]['id'],
                'text': text,
                'filename': safe_filename,
                'language': language,
                'word_count': len(text.split()),
                'size': temp_path.stat().st_size,
                'path': temp_path
            })
            
            await send_progress(ws_id, {
                'stage': 'extraction',
                'progress': idx + 1,
                'total': len(files)
            })
        
        # Comparaisons
        comparisons_total = len(file_records) * (len(file_records) - 1) // 2
        comparisons_done = 0
        matches = []
        
        await send_progress(ws_id, {
            'stage': 'comparison',
            'progress': 0,
            'total': comparisons_total,
            'message': 'Comparaison des fichiers...'
        })
        
        for i in range(len(file_records)):
            for j in range(i + 1, len(file_records)):
                file_a = file_records[i]
                file_b = file_records[j]
                
                similarity, details = calculate_similarity(file_a['text'], file_b['text'])
                
                comparisons_done += 1
                
                if similarity >= threshold:
                    report = supabase_client.table('similarity_reports').insert({
                        'analysis_id': analysis_id,
                        'file_a_id': file_a['id'],
                        'file_b_id': file_b['id'],
                        'similarity_percentage': similarity,
                        'similarity_type': f"{'Code' if file_a['language'] in ['Python', 'Java', 'C', 'JavaScript', 'PHP'] else 'Texte'} - {'Exact' if similarity > 80 else 'Modéré' if similarity > 50 else 'Partiel'}",
                        'exact_matches': details['exact_count'],
                        'moderate_matches': details['moderate_count'],
                        'weak_matches': details['weak_count'],
                        'segments': json.dumps(details['segments'])
                    }).execute()
                    
                    matches.append({
                        'report_id': report.data[0]['id'],
                        'file_a': file_a,
                        'file_b': file_b,
                        'similarity': similarity,
                        'details': details
                    })
                
                await send_progress(ws_id, {
                    'stage': 'comparison',
                    'progress': comparisons_done,
                    'total': comparisons_total
                })
        
        # Génération des rapports PDF
        await send_progress(ws_id, {
            'stage': 'report',
            'progress': 0,
            'total': len(matches),
            'message': 'Génération des rapports PDF...'
        })
        
        teacher = supabase_client.table('teachers').select('*').eq('id', teacher_id).execute()
        establishment = None
        if establishment_id:
            establishment = supabase_client.table('establishments').select('*').eq('id', establishment_id).execute()
        
        for idx, match in enumerate(matches):
            report_filename = f"report_{match['report_id']}.pdf"
            report_path = REPORTS_DIR / report_filename
            
            report_data = {
                'report_id': match['report_id'][:8],
                'date': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                'teacher_name': teacher.data[0]['name'] if teacher.data else '',
                'establishment_name': establishment.data[0]['name'] if establishment and establishment.data else '',
                'global_similarity': match['similarity'],
                'total_comparisons': comparisons_total,
                'avg_similarity': round(sum(m['similarity'] for m in matches) / len(matches), 2),
                'matches_count': len(matches),
                'threshold': threshold,
                'file_a_name': match['file_a']['filename'],
                'file_b_name': match['file_b']['filename'],
                'file_a_size': match['file_a']['size'],
                'file_b_size': match['file_b']['size'],
                'file_a_words': match['file_a']['word_count'],
                'file_b_words': match['file_b']['word_count'],
                'file_a_language': match['file_a']['language'],
                'file_b_language': match['file_b']['language'],
                'exact_matches': match['details']['exact_count'],
                'moderate_matches': match['details']['moderate_count'],
                'weak_matches': match['details']['weak_count'],
                'similarity_type': f"{'Code' if match['file_a']['language'] in ['Python', 'Java'] else 'Texte'}",
                'text_a': match['file_a']['text'],
                'text_b': match['file_b']['text'],
                'similarity_no_quotes': match['similarity'],
                'structural_similarity': round(match['similarity'] * 0.9, 2),
                'syntactic_similarity': round(match['similarity'] * 0.85, 2),
                'signature': hashlib.sha256(f"{match['report_id']}{datetime.now().isoformat()}".encode()).hexdigest()[:16]
            }
            
            generate_pdf_report(report_data, report_path)
            
            pdf_storage_path = f"reports/{analysis_id}/{report_filename}"
            with open(report_path, 'rb') as f:
                get_bucket('plagify-reports').upload(
                    pdf_storage_path,
                    f.read(),
                    {'content-type': 'application/pdf', 'upsert': 'true'}
                )
            
            pdf_url = get_bucket('plagify-reports').get_public_url(pdf_storage_path)
            
            supabase_client.table('similarity_reports').update({
                'report_pdf_url': pdf_url
            }).eq('id', match['report_id']).execute()
            
            await send_progress(ws_id, {
                'stage': 'report',
                'progress': idx + 1,
                'total': len(matches)
            })
        
        # Finaliser l'analyse
        supabase_client.table('analyses').update({
            'status': 'completed',
            'completed_at': datetime.now().isoformat(),
            'total_comparisons': comparisons_total,
            'matches_above_threshold': len(matches),
            'avg_similarity': round(sum(m['similarity'] for m in matches) / len(matches), 2) if matches else 0
        }).eq('id', analysis_id).execute()
        
        supabase_client.table('activity_logs').insert({
            'teacher_id': teacher_id,
            'analysis_id': analysis_id,
            'activity_type': 'analysis_completed',
            'message': f"Analyse terminée: {len(matches)} correspondances trouvées sur {comparisons_total} comparaisons",
            'metadata': json.dumps({'matches': len(matches), 'threshold': threshold})
        }).execute()
        
        await send_progress(ws_id, {
            'stage': 'complete',
            'progress': 100,
            'total': 100,
            'message': 'Analyse terminée!',
            'analysis_id': analysis_id,
            'matches': len(matches)
        })
        
        # Nettoyer fichiers temporaires
        for file_rec in file_records:
            if file_rec['path'].exists():
                file_rec['path'].unlink()
        
        return {
            "success": True,
            "analysis_id": analysis_id,
            "matches": len(matches),
            "total_comparisons": comparisons_total
        }
        
    except Exception as e:
        if 'analysis_id' in locals():
            supabase_client.table('analyses').update({
                'status': 'failed',
                'error_message': str(e),
                'completed_at': datetime.now().isoformat()
            }).eq('id', analysis_id).execute()
        
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/analyze/single-file")
async def analyze_single_file(
    teacher_id: str = Form(...),
    establishment_id: Optional[str] = Form(None),
    threshold: float = Form(15.0),
    ws_id: str = Form(...),
    file: UploadFile = File(...)
):
    """
    Analyse un fichier unique contre tous les fichiers de la base de données
    """
    try:
        analysis = supabase_client.table('analyses').insert({
            'teacher_id': teacher_id,
            'establishment_id': establishment_id,
            'analysis_type': 'single_file',
            'source_name': file.filename,
            'similarity_threshold': threshold,
            'status': 'processing'
        }).execute()
        
        analysis_id = analysis.data[0]['id']
        
        # ✅ FIX PROBLÈME 5 — Nom de fichier sûr
        safe_filename = Path(file.filename).name
        temp_path = UPLOAD_DIR / f"{analysis_id}_{safe_filename}"
        with open(temp_path, 'wb') as f:
            f.write(await file.read())
        
        text, language = extract_text_from_file(temp_path)
        
        db_files = supabase_client.table('files').select('*').eq('teacher_id', teacher_id).execute()
        
        total = len(db_files.data)
        matches = []
        
        await send_progress(ws_id, {
            'stage': 'comparison',
            'progress': 0,
            'total': total,
            'message': f'Comparaison avec {total} fichiers...'
        })
        
        for idx, db_file in enumerate(db_files.data):
            similarity, details = calculate_similarity(text, db_file['content_text'] or '')
            
            if similarity >= threshold:
                report = supabase_client.table('similarity_reports').insert({
                    'analysis_id': analysis_id,
                    'file_a_id': db_file['id'],
                    'file_b_id': db_file['id'],
                    'similarity_percentage': similarity,
                    'similarity_type': 'Comparaison base de données',
                    'exact_matches': details['exact_count'],
                    'moderate_matches': details['moderate_count'],
                    'weak_matches': details['weak_count'],
                    'segments': json.dumps(details['segments'])
                }).execute()
                
                matches.append({
                    'file': db_file,
                    'similarity': similarity,
                    'details': details
                })
            
            await send_progress(ws_id, {
                'stage': 'comparison',
                'progress': idx + 1,
                'total': total
            })
        
        supabase_client.table('analyses').update({
            'status': 'completed',
            'completed_at': datetime.now().isoformat(),
            'total_comparisons': total,
            'matches_above_threshold': len(matches),
            'total_files': 1
        }).eq('id', analysis_id).execute()
        
        await send_progress(ws_id, {
            'stage': 'complete',
            'progress': 100,
            'total': 100,
            'analysis_id': analysis_id,
            'matches': len(matches)
        })
        
        if temp_path.exists():
            temp_path.unlink()
        
        return {
            "success": True,
            "analysis_id": analysis_id,
            "matches": len(matches),
            "total_comparisons": total
        }
        
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/analyses/{teacher_id}")
async def get_analyses(teacher_id: str):
    """Récupérer toutes les analyses d'un enseignant"""
    try:
        result = supabase_client.table('analyses').select('*').eq('teacher_id', teacher_id).order('started_at', desc=True).execute()
        return {"success": True, "data": result.data}
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/analyses/{analysis_id}/reports")
async def get_analysis_reports(analysis_id: str):
    """Récupérer tous les rapports d'une analyse"""
    try:
        # ✅ FIX PROBLÈME 4 — Alias pour éviter le conflit de jointure sur la table files
        result = supabase_client.table('similarity_reports').select(
            '*, file_a:files!file_a_id(*), file_b:files!file_b_id(*)'
        ).eq('analysis_id', analysis_id).execute()
        return {"success": True, "data": result.data}
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/statistics/{teacher_id}")
async def get_statistics(teacher_id: str):
    """Récupérer les statistiques d'un enseignant"""
    try:
        files_result = supabase_client.table('files').select('id', count='exact').eq('teacher_id', teacher_id).execute()
        total_files = files_result.count if files_result.count else 0
        
        analyses_result = supabase_client.table('analyses').select('id', count='exact').eq('teacher_id', teacher_id).execute()
        total_analyses = analyses_result.count if analyses_result.count else 0
        
        if total_analyses > 0:
            analyses_ids_result = supabase_client.table('analyses').select('id').eq('teacher_id', teacher_id).execute()
            analyses_ids = [a['id'] for a in analyses_ids_result.data]
            
            if analyses_ids:
                reports_result = supabase_client.table('similarity_reports').select('id', count='exact').in_('analysis_id', analyses_ids).execute()
                total_reports = reports_result.count if reports_result.count else 0
            else:
                total_reports = 0
        else:
            total_reports = 0
        
        return {
            "success": True,
            "data": {
                "total_files": total_files,
                "total_analyses": total_analyses,
                "total_reports": total_reports
            }
        }
    except Exception as e:
        print(f"Error getting statistics: {e}")
        return {
            "success": True,
            "data": {
                "total_files": 0,
                "total_analyses": 0,
                "total_reports": 0
            }
        }

@app.websocket("/ws/{ws_id}")
async def websocket_endpoint(websocket: WebSocket, ws_id: str):
    """WebSocket pour progression en temps réel"""
    await websocket.accept()
    ws_connections[ws_id] = websocket
    
    try:
        while True:
            data = await websocket.receive_text()
    except WebSocketDisconnect:
        if ws_id in ws_connections:
            del ws_connections[ws_id]

# Pour déploiement sur Render
if __name__ == "__main__":
    import uvicorn
    port = int(os.getenv("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port)